"""
Test that generated documents (PDF, XLSX, PPTX, DOCX) are persisted to PostgreSQL.

These tests verify that when documents are created in Pyodide, they are
properly saved to the sandbox_filesystem table in PostgreSQL.
"""

import os
import sys

import asyncpg
import pytest
from dotenv import load_dotenv

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))
load_dotenv()

from mayflower_sandbox.sandbox_executor import SandboxExecutor  # noqa: E402


@pytest.fixture
async def db_pool():
    """Create test database connection pool."""
    db_config = {
        "host": os.getenv("POSTGRES_HOST", "localhost"),
        "database": os.getenv("POSTGRES_DB", "mayflower_test"),
        "user": os.getenv("POSTGRES_USER", "postgres"),
        "password": os.getenv("POSTGRES_PASSWORD", "postgres"),
        "port": int(os.getenv("POSTGRES_PORT", "5432")),
    }

    pool = await asyncpg.create_pool(**db_config)

    # Ensure session exists
    async with pool.acquire() as conn:
        await conn.execute(
            """
            INSERT INTO sandbox_sessions (thread_id, expires_at)
            VALUES ('doc_persist_test', NOW() + INTERVAL '1 day')
            ON CONFLICT (thread_id) DO NOTHING
        """
        )

    yield pool
    await pool.close()


@pytest.fixture
async def executor(db_pool):
    """Create sandbox executor for tests."""
    return SandboxExecutor(db_pool, "doc_persist_test", allow_net=True)


@pytest.fixture
async def clean_files(db_pool):
    """Clean files before each test."""
    async with db_pool.acquire() as conn:
        await conn.execute("DELETE FROM sandbox_filesystem WHERE thread_id = 'doc_persist_test'")
    yield


async def test_excel_persistence(executor, db_pool, clean_files):
    """Test Excel file is persisted to PostgreSQL."""
    code = """
import micropip
await micropip.install('openpyxl')

import openpyxl
from openpyxl.styles import Font

# Create workbook
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Sales Data"

# Add headers
ws['A1'] = 'Product'
ws['B1'] = 'Quantity'
ws['C1'] = 'Price'

# Style headers
for cell in ['A1', 'B1', 'C1']:
    ws[cell].font = Font(bold=True)

# Add data
ws['A2'] = 'Laptop'
ws['B2'] = 5
ws['C2'] = 1200

# Save to filesystem
wb.save('/tmp/sales.xlsx')
print("Excel file created")
"""

    result = await executor.execute(code)

    assert result.success is True, f"Execution failed: {result.stderr}"
    assert "Excel file created" in result.stdout
    assert result.created_files is not None
    assert "/tmp/sales.xlsx" in result.created_files

    # Verify file was saved to PostgreSQL
    async with db_pool.acquire() as conn:
        file_data = await conn.fetchrow(
            """
            SELECT content, content_type, size
            FROM sandbox_filesystem
            WHERE thread_id = 'doc_persist_test' AND file_path = '/tmp/sales.xlsx'
        """
        )

        assert file_data is not None, "Excel file not found in PostgreSQL"
        assert len(file_data["content"]) > 0, "Excel file content is empty"
        assert file_data["size"] > 0, "Excel file size is 0"
        # XLSX files are ZIP archives, should start with PK
        assert file_data["content"][:2] == b"PK", "Excel file has invalid format"


async def test_pdf_persistence(executor, db_pool, clean_files):
    """Test PDF file is persisted to PostgreSQL."""
    code = """
import micropip
await micropip.install('fpdf2')

from fpdf import FPDF

# Create PDF
pdf = FPDF()
pdf.add_page()
pdf.set_font("helvetica", size=12)
pdf.cell(200, 10, txt="Q4 Financial Report", ln=True, align='C')
pdf.cell(200, 10, txt="Revenue increased 25%", ln=True)

# Save to filesystem
pdf.output('/tmp/report.pdf')
print("PDF file created")
"""

    result = await executor.execute(code)

    assert result.success is True, f"Execution failed: {result.stderr}"
    assert "PDF file created" in result.stdout
    assert result.created_files is not None
    assert "/tmp/report.pdf" in result.created_files

    # Verify file was saved to PostgreSQL
    async with db_pool.acquire() as conn:
        file_data = await conn.fetchrow(
            """
            SELECT content, content_type, size
            FROM sandbox_filesystem
            WHERE thread_id = 'doc_persist_test' AND file_path = '/tmp/report.pdf'
        """
        )

        assert file_data is not None, "PDF file not found in PostgreSQL"
        assert len(file_data["content"]) > 0, "PDF file content is empty"
        assert file_data["size"] > 0, "PDF file size is 0"
        assert file_data["content_type"] == "application/pdf"
        # PDF files start with %PDF
        assert file_data["content"][:4] == b"%PDF", "PDF file has invalid format"


async def test_powerpoint_persistence(executor, db_pool, clean_files):
    """Test PowerPoint file is persisted to PostgreSQL."""
    code = """
import micropip
await micropip.install('python-pptx')

from pptx import Presentation

# Create presentation
prs = Presentation()

# Add title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Company Overview 2024"
subtitle.text = "Annual Review"

# Add content slide
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Key Achievements'
tf = body_shape.text_frame
tf.text = 'Revenue growth 30%'

# Save to filesystem
prs.save('/tmp/presentation.pptx')
print("PowerPoint file created")
"""

    result = await executor.execute(code)

    assert result.success is True, f"Execution failed: {result.stderr}"
    assert "PowerPoint file created" in result.stdout
    assert result.created_files is not None
    assert "/tmp/presentation.pptx" in result.created_files

    # Verify file was saved to PostgreSQL
    async with db_pool.acquire() as conn:
        file_data = await conn.fetchrow(
            """
            SELECT content, content_type, size
            FROM sandbox_filesystem
            WHERE thread_id = 'doc_persist_test' AND file_path = '/tmp/presentation.pptx'
        """
        )

        assert file_data is not None, "PowerPoint file not found in PostgreSQL"
        assert len(file_data["content"]) > 0, "PowerPoint file content is empty"
        assert file_data["size"] > 0, "PowerPoint file size is 0"
        # PPTX files are ZIP archives, should start with PK
        assert file_data["content"][:2] == b"PK", "PowerPoint file has invalid format"


async def test_word_persistence(executor, db_pool, clean_files):
    """Test Word document is persisted to PostgreSQL."""
    code = """
import micropip
await micropip.install('python-docx')

from docx import Document

# Create document
doc = Document()
doc.add_heading('Project Status Report', 0)

doc.add_heading('Overview', level=1)
doc.add_paragraph('Project is on track for Q4 delivery.')

doc.add_heading('Key Achievements', level=1)
doc.add_paragraph('Completed phase 1', style='List Bullet')
doc.add_paragraph('Launched beta version', style='List Bullet')

# Save to filesystem
doc.save('/tmp/report.docx')
print("Word document created")
"""

    result = await executor.execute(code)

    assert result.success is True, f"Execution failed: {result.stderr}"
    assert "Word document created" in result.stdout
    assert result.created_files is not None
    assert "/tmp/report.docx" in result.created_files

    # Verify file was saved to PostgreSQL
    async with db_pool.acquire() as conn:
        file_data = await conn.fetchrow(
            """
            SELECT content, content_type, size
            FROM sandbox_filesystem
            WHERE thread_id = 'doc_persist_test' AND file_path = '/tmp/report.docx'
        """
        )

        assert file_data is not None, "Word document not found in PostgreSQL"
        assert len(file_data["content"]) > 0, "Word document content is empty"
        assert file_data["size"] > 0, "Word document size is 0"
        # DOCX files are ZIP archives, should start with PK
        assert file_data["content"][:2] == b"PK", "Word document has invalid format"


async def test_multiple_documents_persistence(executor, db_pool, clean_files):
    """Test multiple documents created in same execution are all persisted."""
    code = """
import micropip
await micropip.install('fpdf2')
await micropip.install('openpyxl')

from fpdf import FPDF
import openpyxl

# Create PDF
pdf = FPDF()
pdf.add_page()
pdf.set_font("helvetica", size=12)
pdf.cell(200, 10, txt="Sales Report", ln=True)
pdf.output('/tmp/sales_report.pdf')

# Create Excel
wb = openpyxl.Workbook()
ws = wb.active
ws['A1'] = 'Item'
ws['B1'] = 'Value'
ws['A2'] = 'Revenue'
ws['B2'] = 100000
wb.save('/tmp/sales_data.xlsx')

print("Both files created")
"""

    result = await executor.execute(code)

    assert result.success is True, f"Execution failed: {result.stderr}"
    assert "Both files created" in result.stdout
    assert result.created_files is not None
    assert "/tmp/sales_report.pdf" in result.created_files
    assert "/tmp/sales_data.xlsx" in result.created_files

    # Verify both files are in PostgreSQL
    async with db_pool.acquire() as conn:
        pdf_file = await conn.fetchrow(
            """
            SELECT content FROM sandbox_filesystem
            WHERE thread_id = 'doc_persist_test' AND file_path = '/tmp/sales_report.pdf'
        """
        )
        xlsx_file = await conn.fetchrow(
            """
            SELECT content FROM sandbox_filesystem
            WHERE thread_id = 'doc_persist_test' AND file_path = '/tmp/sales_data.xlsx'
        """
        )

        assert pdf_file is not None, "PDF not found in PostgreSQL"
        assert xlsx_file is not None, "Excel not found in PostgreSQL"
        assert pdf_file["content"][:4] == b"%PDF", "PDF has invalid format"
        assert xlsx_file["content"][:2] == b"PK", "Excel has invalid format"
